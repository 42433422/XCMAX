"""办公媒体（PPTX / PDF）→ 模版库 analyze 兼容载荷。

同步路径：文本/占位符优先；扫描件 PDF 走 OCR；可选 VLM 识图增强。
不强制拉起 full-read 员工 Mod（避免 analyze HTTP 依赖异步员工运行时）。
"""

from __future__ import annotations

import logging
import os
import re
import zipfile
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from app.utils.operational_errors import RECOVERABLE_ERRORS

logger = logging.getLogger(__name__)

_PLACEHOLDER_PATTERNS = (
    re.compile(r"\{\{\s*([^}]+?)\s*\}\}"),
    re.compile(r"\{\%\s*([^\%]+?)\s*\%\}"),
    re.compile(r"\$\{\s*([^}]+?)\s*\}"),
    re.compile(r"\[\[\s*([^\]]+?)\s*\]\]"),
)

_A_T = "{http://schemas.openxmlformats.org/drawingml/2006/main}t"
_W_T = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"


def extract_placeholder_tokens(text: str) -> list[str]:
    raw_tokens: list[str] = []
    for pat in _PLACEHOLDER_PATTERNS:
        for m in pat.finditer(text or ""):
            token = str(m.group(1) or "").strip()
            if token and token not in raw_tokens:
                raw_tokens.append(token)
    return raw_tokens


def fields_from_tokens(tokens: list[str]) -> list[dict[str, Any]]:
    return [{"label": t, "value": "", "type": "dynamic"} for t in tokens]


def fields_from_text_lines(text: str, *, limit: int = 40) -> list[dict[str, Any]]:
    """无占位符时，用非空行生成可编辑预览字段（便于入库后人工标注）。"""
    fields: list[dict[str, Any]] = []
    seen: set[str] = set()
    for line in (text or "").splitlines():
        label = " ".join(line.strip().split())
        if not label or label in seen:
            continue
        seen.add(label)
        fields.append({"label": label[:80], "value": "", "type": "dynamic"})
        if len(fields) >= limit:
            break
    return fields


def _collect_ooxml_text(xml_bytes: bytes, tag: str) -> str:
    try:
        root = ET.fromstring(xml_bytes)
    except ET.ParseError:
        return ""
    parts: list[str] = []
    for node in root.iter(tag):
        if node.text:
            parts.append(node.text)
        if node.tail:
            parts.append(node.tail)
    return "".join(parts)


def _extract_pptx_text_via_zip(file_path: str) -> tuple[str, int]:
    blobs: list[str] = []
    slide_count = 0
    with zipfile.ZipFile(file_path, "r") as zf:
        names = sorted(
            n
            for n in zf.namelist()
            if re.match(r"^ppt/slides/slide\d+\.xml$", n)
            or re.match(r"^ppt/notesSlides/notesSlide\d+\.xml$", n)
        )
        for name in names:
            if "/slides/slide" in name:
                slide_count += 1
            try:
                blobs.append(_collect_ooxml_text(zf.read(name), _A_T))
            except (KeyError, OSError, ET.ParseError) as exc:
                logger.debug("skip pptx part %s: %s", name, exc)
    return "\n".join(b for b in blobs if b), slide_count


def _extract_pptx_text_via_lib(file_path: str) -> tuple[str, int]:
    from pptx import Presentation

    prs = Presentation(file_path)
    chunks: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        chunks.append(f"[slide {idx}]")
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            try:
                for para in shape.text_frame.paragraphs:
                    t = (para.text or "").strip()
                    if t:
                        chunks.append(t)
            except RECOVERABLE_ERRORS:
                continue
        if getattr(slide, "has_notes_slide", False):
            try:
                notes = slide.notes_slide.notes_text_frame.text
                if notes and notes.strip():
                    chunks.append(notes.strip())
            except RECOVERABLE_ERRORS:
                pass
    return "\n".join(chunks), len(prs.slides)


def extract_pptx_document_text(file_path: str) -> dict[str, Any]:
    engine = "zip_ooxml"
    try:
        text, slide_count = _extract_pptx_text_via_lib(file_path)
        engine = "python-pptx"
    except (RECOVERABLE_ERRORS, ImportError, ModuleNotFoundError):
        logger.debug("python-pptx unavailable, fallback zip", exc_info=True)
        text, slide_count = _extract_pptx_text_via_zip(file_path)
    return {
        "text": text,
        "slide_count": slide_count,
        "engine": engine,
        "char_count": len(text or ""),
    }


def extract_pdf_document_text(file_path: str) -> dict[str, Any]:
    from pypdf import PdfReader

    reader = PdfReader(file_path)
    pages: list[str] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            page_text = page.extract_text() or ""
        except RECOVERABLE_ERRORS:
            page_text = ""
        if page_text.strip():
            pages.append(f"[page {i}]\n{page_text}")
    text = "\n\n".join(pages)
    return {
        "text": text,
        "page_count": len(reader.pages),
        "engine": "pypdf",
        "char_count": len(text),
    }


def _ocr_pdf_plaintext(file_path: str) -> dict[str, Any]:
    from app.application.shipment_excel_etl_ocr import ocr_source_to_workbook

    result = ocr_source_to_workbook(file_path)
    if not result.get("success"):
        return {
            "success": False,
            "text": "",
            "message": str(result.get("message") or "OCR 失败"),
            "engine": "ocr",
        }
    # OCR 产出临时 xlsx；读回网格拼文本，供占位符/字段提取
    grid_path = str(result.get("workbook_path") or "")
    lines: list[str] = []
    try:
        from openpyxl import load_workbook

        wb = load_workbook(grid_path, read_only=True, data_only=True)
        ws = wb.active
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                lines.append("\t".join(cells))
        wb.close()
    except RECOVERABLE_ERRORS as exc:
        return {"success": False, "text": "", "message": f"OCR 结果读取失败: {exc}", "engine": "ocr"}
    text = "\n".join(lines)
    return {
        "success": True,
        "text": text,
        "engine": "ocr",
        "char_count": len(text),
        "workbook_path": grid_path,
        "message": "OCR 成功",
    }


def _vlm_enrich_enabled() -> bool:
    return os.environ.get("FHD_TEMPLATE_VLM_ENRICH", "").strip().lower() in {
        "1",
        "true",
        "yes",
        "on",
    }


def _sync_vlm_describe_first_image(image_bytes: bytes, *, hint: str) -> dict[str, Any] | None:
    """可选：用平台 VLM 路由描述一张图（失败则静默跳过）。"""
    if not _vlm_enrich_enabled() or not image_bytes:
        return None
    try:
        import asyncio
        import base64

        from app.infrastructure.llm.vlm_route import resolve_vlm_route
        from app.mod_sdk.mod_employee_llm import mod_employee_complete

        route = resolve_vlm_route()
        if not route.get("ok"):
            return {"vlm_ok": False, "message": route.get("message"), "route": route}

        b64 = base64.b64encode(image_bytes).decode("ascii")
        url = f"data:image/png;base64,{b64}"
        messages = [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": (
                            f"{hint} 请用中文提取可见文字与关键字段，输出 JSON："
                            '{"description":"","detected_text":"","fields":[{"label":"","value":""}]}'
                        ),
                    },
                    {"type": "image_url", "image_url": {"url": url}},
                ],
            }
        ]

        async def _run() -> dict[str, Any]:
            return await mod_employee_complete(messages, max_tokens=800, temperature=0.1)

        try:
            loop = asyncio.get_running_loop()
        except RuntimeError:
            loop = None
        if loop and loop.is_running():
            # analyze 路径通常无事件循环；若已有循环则跳过以免死锁
            return {"vlm_ok": False, "message": "event_loop_busy_skip_vlm", "route": route}

        result = asyncio.run(_run())
        if not result.get("success"):
            return {
                "vlm_ok": False,
                "message": result.get("error") or "vlm_failed",
                "route": route,
            }
        content = str(result.get("content") or "").strip()
        return {
            "vlm_ok": True,
            "description": content[:2000],
            "route": {"provider": route.get("provider"), "model": route.get("model")},
        }
    except RECOVERABLE_ERRORS as exc:
        logger.debug("vlm enrich skipped: %s", exc, exc_info=True)
        return {"vlm_ok": False, "message": str(exc)}


def _maybe_pdf_page_image_bytes(file_path: str) -> bytes | None:
    try:
        from app.application.shipment_excel_etl_ocr import _load_image_arrays
        from PIL import Image
        import io

        images = _load_image_arrays(Path(file_path))
        if not images:
            return None
        img = Image.fromarray(images[0])
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()
    except RECOVERABLE_ERRORS:
        return None


def build_pptx_template_analysis(
    file_path: str,
    *,
    template_name: str = "",
    original_filename: str = "",
) -> dict[str, Any]:
    meta = extract_pptx_document_text(file_path)
    text = str(meta.get("text") or "")
    tokens = extract_placeholder_tokens(text)
    fields = fields_from_tokens(tokens) if tokens else fields_from_text_lines(text)
    if not fields:
        return {
            "success": False,
            "message": "未能从 PPTX 提取占位符或正文。请确认文件含可提取文本。",
            "engine": meta.get("engine"),
        }
    name = template_name or Path(original_filename or file_path).stem
    snippet = " ".join(text.split())
    if len(snippet) > 400:
        snippet = snippet[:400] + "…"
    return {
        "success": True,
        "template_name": name,
        "template_type": "pptx",
        "fields": fields,
        "preview_data": {
            "file_path": file_path,
            "original_filename": original_filename or Path(file_path).name,
            "placeholders": tokens,
            "text_snippet": snippet,
            "slide_count": meta.get("slide_count"),
            "engine": meta.get("engine"),
            "parser": "office_template_media_bridge",
        },
    }


def build_pdf_template_analysis(
    file_path: str,
    *,
    template_name: str = "",
    original_filename: str = "",
) -> dict[str, Any]:
    warnings: list[str] = []
    try:
        meta = extract_pdf_document_text(file_path)
    except RECOVERABLE_ERRORS as exc:
        return {"success": False, "message": f"PDF 文本提取失败：{exc}"}

    text = str(meta.get("text") or "")
    engine = str(meta.get("engine") or "pypdf")
    ocr_meta: dict[str, Any] = {}
    # 扫描件：文本过少则 OCR
    if len(text.strip()) < 40:
        ocr = _ocr_pdf_plaintext(file_path)
        if ocr.get("success") and ocr.get("text"):
            text = str(ocr["text"])
            engine = "ocr"
            ocr_meta = {
                "workbook_path": ocr.get("workbook_path"),
                "char_count": ocr.get("char_count"),
            }
        else:
            warnings.append(str(ocr.get("message") or "OCR 未产出文本"))

    vlm_sidecar = None
    if len(text.strip()) < 40 or _vlm_enrich_enabled():
        img_bytes = _maybe_pdf_page_image_bytes(file_path)
        if img_bytes:
            vlm_sidecar = _sync_vlm_describe_first_image(
                img_bytes, hint="这是 PDF 首页渲染图。"
            )
            if vlm_sidecar and vlm_sidecar.get("vlm_ok") and vlm_sidecar.get("description"):
                # 把 VLM 描述并入文本，便于占位符/字段生成
                text = f"{text}\n{vlm_sidecar['description']}".strip()

    tokens = extract_placeholder_tokens(text)
    fields = fields_from_tokens(tokens) if tokens else fields_from_text_lines(text)
    if not fields:
        return {
            "success": False,
            "message": "未能从 PDF 提取可用文本（含 OCR/VLM）。请确认文件可读或启用 VLM/OCR。",
            "warnings": warnings,
            "engine": engine,
            "vlm": vlm_sidecar,
        }

    name = template_name or Path(original_filename or file_path).stem
    snippet = " ".join(text.split())
    if len(snippet) > 400:
        snippet = snippet[:400] + "…"
    preview: dict[str, Any] = {
        "file_path": file_path,
        "original_filename": original_filename or Path(file_path).name,
        "placeholders": tokens,
        "text_snippet": snippet,
        "page_count": meta.get("page_count"),
        "engine": engine,
        "parser": "office_template_media_bridge",
        "warnings": warnings,
    }
    if ocr_meta:
        preview["ocr"] = ocr_meta
    if vlm_sidecar:
        preview["vlm"] = vlm_sidecar
    return {
        "success": True,
        "template_name": name,
        "template_type": "pdf",
        "fields": fields,
        "preview_data": preview,
    }
