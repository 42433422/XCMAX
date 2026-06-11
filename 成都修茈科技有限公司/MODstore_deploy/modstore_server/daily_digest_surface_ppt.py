"""每日摘要 · 三端页面巡检截图 → PowerPoint（.pptx）。

把 :mod:`modstore_server.daily_digest_surface_audit` 抓到的 P-W / P-S / P-App
关键页截图与「对应员工 AI 分析」拼成一份 PPT：

- 封面页：日期 + 三端通过 / 异常概览；
- 每条产线一页分线分析（AI 分析 + 对应员工）；
- 每张截图一页：左侧整页截图（保持比例缩放），右侧标题 / URL / HTTP / console / 分析。

随每日摘要邮件作为附件投递（见 :func:`modstore_server.email_service.send_html_email_with_attachments`）。

环境变量：
- ``MODSTORE_DAILY_SURFACE_PPT_ENABLED``（默认 ``1``）：设为 ``0`` 不生成 PPT。
- ``MODSTORE_DAILY_SURFACE_PPT_SAVE_DIR``：PPT 输出目录（默认 ``playwright-report/digest-surfaces`` 相对仓库根，与截图同级）。
"""

from __future__ import annotations

import logging
import os
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)

# 16:9 宽屏（EMU 由 python-pptx 的 Inches 转换）
_SLIDE_W_IN = 13.333
_SLIDE_H_IN = 7.5

_LANE_TITLES = {
    "P-W": "P-W · 获客网站（xiu-ci.com）",
    "P-S": "P-S · MODstore 软件面（market）",
    "P-App": "P-App · 移动端 / adb 原生屏",
}
_LANE_ORDER = ["P-W", "P-S", "P-App"]


def _enabled() -> bool:
    raw = (os.environ.get("MODSTORE_DAILY_SURFACE_PPT_ENABLED", "1") or "").strip().lower()
    return raw not in ("0", "false", "no", "off")


def _repo_root() -> Path:
    try:
        from modstore_server.daily_digest import _repo_root as root_fn

        return Path(root_fn())
    except Exception:  # noqa: BLE001
        return Path(os.environ.get("MODSTORE_REPO_ROOT", ".")).resolve()


def _save_dir(day: str) -> Path:
    raw = (
        os.environ.get("MODSTORE_DAILY_SURFACE_PPT_SAVE_DIR") or "playwright-report/digest-surfaces"
    ).strip()
    out = _repo_root() / raw / day
    out.mkdir(parents=True, exist_ok=True)
    return out


def _fit_image_box(
    img_path: Path,
    *,
    max_w_in: float,
    max_h_in: float,
) -> Tuple[float, float]:
    """按图片真实宽高比，等比缩放到 ``max_w x max_h`` 盒子内，返回 (w_in, h_in)。"""
    try:
        from PIL import Image

        with Image.open(img_path) as im:
            w_px, h_px = im.size
        if w_px <= 0 or h_px <= 0:
            return max_w_in, max_h_in
        ratio = min(max_w_in / w_px, max_h_in / h_px)
        return max(0.5, w_px * ratio), max(0.5, h_px * ratio)
    except Exception:  # noqa: BLE001
        return max_w_in, max_h_in


def _lane_analysis_for(report: Dict[str, Any], lane: str) -> Dict[str, Any]:
    la = report.get("lane_analysis") if isinstance(report.get("lane_analysis"), dict) else {}
    row = la.get(lane) if isinstance(la, dict) else None
    return row if isinstance(row, dict) else {}


def build_surface_audit_pptx(
    report: Dict[str, Any],
    *,
    out_path: Optional[Path] = None,
) -> Dict[str, Any]:
    """把巡检 report 拼成 .pptx。

    返回 ``{ok, path, slides, lanes, error, skipped}``：
    - python-pptx 未安装 → ``ok=False`` + ``error``；
    - 无截图（results 为空或都没存盘）→ ``skipped=True``。
    """
    if not _enabled():
        return {"ok": True, "skipped": True, "reason": "disabled", "path": "", "slides": 0}

    results = report.get("results") if isinstance(report.get("results"), list) else []
    if not results:
        return {"ok": True, "skipped": True, "reason": "no results", "path": "", "slides": 0}

    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        return {"ok": False, "error": f"未安装 python-pptx：{exc}", "path": "", "slides": 0}

    day = str(report.get("day") or datetime.now(timezone.utc).strftime("%Y-%m-%d"))
    if out_path is None:
        out_path = _save_dir(day) / f"surface-audit-{day}.pptx"
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_W_IN)
    prs.slide_height = Inches(_SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    navy = RGBColor(0x1E, 0x3A, 0x8A)
    slate = RGBColor(0x1E, 0x29, 0x3B)
    grey = RGBColor(0x64, 0x74, 0x8B)
    red = RGBColor(0xB9, 0x1C, 0x1C)
    green = RGBColor(0x04, 0x78, 0x57)
    indigo = RGBColor(0x43, 0x38, 0xCA)

    def _add_text(slide, left, top, width, height, lines: List[Tuple[str, int, Any, bool]]):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for text, size, color, bold in lines:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = p.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
        return box

    lane_counts: Dict[str, Dict[str, int]] = {}
    for r in results:
        lane = str(r.get("lane") or "?")
        c = lane_counts.setdefault(lane, {"ok": 0, "bad": 0})
        if (r.get("status") or 0) < 400 and not r.get("error"):
            c["ok"] += 1
        else:
            c["bad"] += 1

    # ── 封面 ─────────────────────────────────────────────────────────────
    cover = prs.slides.add_slide(blank)
    cover.shapes.add_textbox(Inches(0), Inches(0), Inches(_SLIDE_W_IN), Inches(_SLIDE_H_IN))
    _add_text(
        cover,
        0.8,
        2.0,
        _SLIDE_W_IN - 1.6,
        1.2,
        [("MODstore 三端页面巡检", 40, navy, True)],
    )
    overview = "  ·  ".join(
        f"{lane} 正常 {lane_counts.get(lane, {}).get('ok', 0)} / 异常 {lane_counts.get(lane, {}).get('bad', 0)}"
        for lane in _LANE_ORDER
        if lane in lane_counts
    )
    _add_text(
        cover,
        0.8,
        3.3,
        _SLIDE_W_IN - 1.6,
        2.0,
        [
            (f"日期 {day}", 18, grey, False),
            (overview or "无巡检数据", 16, slate, False),
            (
                "P-W 网站截图+分析 · P-S 软件截图+分析 · P-App adb 模拟器原生屏截图+分析",
                13,
                grey,
                False,
            ),
        ],
    )

    slide_count = 1

    # ── 分线分析页 + 截图页 ──────────────────────────────────────────────
    for lane in _LANE_ORDER:
        lane_rows = [r for r in results if str(r.get("lane")) == lane]
        if not lane_rows:
            continue
        la = _lane_analysis_for(report, lane)
        analysis_md = str(la.get("markdown") or "").strip()
        owners = la.get("owners") or []

        # 分线分析页
        sec = prs.slides.add_slide(blank)
        _add_text(
            sec, 0.6, 0.4, _SLIDE_W_IN - 1.2, 0.9, [(_LANE_TITLES.get(lane, lane), 28, navy, True)]
        )
        owner_line = (
            f"对应员工：{', '.join(str(o) for o in owners[:6])}" if owners else "对应员工：—"
        )
        lines: List[Tuple[str, int, Any, bool]] = [
            (owner_line, 14, indigo, True),
            ("", 8, grey, False),
        ]
        if analysis_md:
            for ln in analysis_md.splitlines():
                if ln.strip():
                    lines.append(
                        (ln.strip(), 15, slate, ln.strip().startswith(("现状", "异常", "改进建议")))
                    )
        else:
            lines.append(("（本产线暂无 AI 分析）", 14, grey, False))
        _add_text(sec, 0.6, 1.4, _SLIDE_W_IN - 1.2, 5.5, lines)
        slide_count += 1

        # 每张截图一页
        for r in lane_rows:
            saved = str(r.get("screenshot_saved") or "").strip()
            slide = prs.slides.add_slide(blank)
            st = r.get("status")
            bad = (st or 500) >= 400 or r.get("error")
            status_color = red if bad else green
            _add_text(
                slide,
                0.5,
                0.3,
                _SLIDE_W_IN - 1.0,
                0.8,
                [(f"{lane} · {r.get('name') or ''}", 22, navy, True)],
            )
            _add_text(
                slide,
                0.5,
                1.05,
                _SLIDE_W_IN - 1.0,
                0.5,
                [(f"HTTP {st or '—'} · {r.get('viewport') or ''}", 13, status_color, True)],
            )

            img_left, img_top = 0.5, 1.6
            img_max_w, img_max_h = 6.2, 5.4
            if saved and Path(saved).is_file():
                w_in, h_in = _fit_image_box(Path(saved), max_w_in=img_max_w, max_h_in=img_max_h)
                try:
                    slide.shapes.add_picture(
                        saved,
                        Inches(img_left),
                        Inches(img_top),
                        width=Inches(w_in),
                        height=Inches(h_in),
                    )
                except Exception:  # noqa: BLE001
                    logger.warning("surface ppt: add_picture failed path=%s", saved)
                    _add_text(
                        slide,
                        img_left,
                        img_top,
                        img_max_w,
                        1.0,
                        [("（截图加载失败）", 14, grey, False)],
                    )
            else:
                _add_text(
                    slide,
                    img_left,
                    img_top,
                    img_max_w,
                    1.0,
                    [("（本页未保存截图）", 14, grey, False)],
                )

            # 右栏：元信息 + 分析
            info_left = 7.0
            info_w = _SLIDE_W_IN - info_left - 0.4
            info_lines: List[Tuple[str, int, Any, bool]] = [
                ("URL", 11, grey, True),
                (str(r.get("url") or ""), 12, slate, False),
                ("页面标题", 11, grey, True),
                (str(r.get("title") or "—"), 12, slate, False),
            ]
            if r.get("error"):
                info_lines += [
                    ("抓取错误", 11, red, True),
                    (str(r.get("error"))[:300], 11, red, False),
                ]
            ce = r.get("console_errors") or []
            if ce:
                info_lines.append(
                    (f"console 报错（{len(ce)}）", 11, RGBColor(0xB4, 0x53, 0x09), True)
                )
                for x in ce[:4]:
                    info_lines.append((f"· {str(x)[:160]}", 10, RGBColor(0x92, 0x40, 0x0E), False))
            analysis = str(r.get("analysis") or "").strip()
            if analysis:
                info_lines.append(("AI 分析", 12, indigo, True))
                for ln in analysis.splitlines():
                    if ln.strip():
                        info_lines.append((ln.strip(), 11, slate, False))
            _add_text(slide, info_left, 1.6, info_w, 5.4, info_lines)
            slide_count += 1

    try:
        prs.save(str(out_path))
    except Exception as exc:  # noqa: BLE001
        logger.exception("surface ppt: save failed")
        return {"ok": False, "error": f"保存 pptx 失败：{exc}", "path": "", "slides": slide_count}

    return {
        "ok": True,
        "skipped": False,
        "path": str(out_path),
        "slides": slide_count,
        "lanes": [lane for lane in _LANE_ORDER if any(str(r.get("lane")) == lane for r in results)],
        "day": day,
    }
