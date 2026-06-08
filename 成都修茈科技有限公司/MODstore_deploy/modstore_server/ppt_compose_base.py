"""Build base PPTX from edit plan slides (compose-first)."""

from __future__ import annotations

import io
import shutil
from pathlib import Path
from typing import Any, Dict, List, Optional

EMU_PER_INCH = 914400


def _resolve_image_path(workspace: Path, ref: str) -> Optional[Path]:
    ref = (ref or "").strip().replace("\\", "/")
    if not ref:
        return None
    candidates = [
        workspace / ref,
        workspace / "outputs" / ref,
        workspace / "outputs" / "images" / ref,
    ]
    if "/" in ref:
        candidates.insert(0, workspace / "outputs" / ref)
    for cand in candidates:
        if cand.is_file():
            return cand
    # category subdirs
    images_root = workspace / "outputs" / "images"
    if images_root.is_dir():
        for sub in images_root.rglob(Path(ref).name):
            if sub.is_file():
                return sub
    return None


def create_compose_deck(
    plan: Dict[str, Any],
    output_path: Path,
    *,
    workspace_root: Optional[Path] = None,
) -> Dict[str, Any]:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError("服务器未安装 python-pptx，暂不能生成 PPT 文件") from exc

    workspace = Path(workspace_root or output_path.parent)
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    title = str(plan.get("title") or "演示文稿")[:120]
    slides_in = plan.get("slides") if isinstance(plan.get("slides"), list) else []
    if not slides_in:
        slides_in = [{"index": 1, "title": title, "bullets": ["暂无内容"], "images": []}]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    logical_map: Dict[str, Dict[str, Any]] = {}
    media_count = 0

    # Cover
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    bg = cover.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(10, 12, 20)
    box = cover.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.6), Inches(1.4))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(248, 250, 252)

    for idx, raw in enumerate(slides_in[:30]):
        if not isinstance(raw, dict):
            continue
        slide_title = str(raw.get("title") or f"第 {idx + 1} 页")[:120]
        bullets = raw.get("bullets") if isinstance(raw.get("bullets"), list) else []
        notes = str(raw.get("notes") or "")[:8000]

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(15, 23, 42)

        badge = slide.shapes.add_textbox(Inches(0.75), Inches(0.42), Inches(1.0), Inches(0.3))
        badge.text_frame.text = f"{idx + 1:02d}"
        badge.text_frame.paragraphs[0].font.size = Pt(11)
        badge.text_frame.paragraphs[0].font.color.rgb = RGBColor(125, 211, 252)

        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.82), Inches(11.8), Inches(0.72))
        title_box.text_frame.text = slide_title
        tp = title_box.text_frame.paragraphs[0]
        tp.font.size = Pt(28)
        tp.font.bold = True
        tp.font.color.rgb = RGBColor(248, 250, 252)

        body = slide.shapes.add_textbox(Inches(1.0), Inches(1.95), Inches(11.1), Inches(3.2))
        tf = body.text_frame
        tf.word_wrap = True
        tf.clear()
        for i, line in enumerate(bullets[:7]):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = str(line)[:180]
            para.font.size = Pt(18)
            para.font.color.rgb = RGBColor(226, 232, 240)

        images = raw.get("images") if isinstance(raw.get("images"), list) else []
        pic_left = 0.6
        pic_count = min(len(images), 4)
        pic_width = 2.4 if pic_count else 0
        for j, img in enumerate(images[:4]):
            if not isinstance(img, dict):
                continue
            ref = str(img.get("path") or img.get("image_ref") or "").strip()
            img_path = _resolve_image_path(workspace, ref)
            if not img_path or not img_path.is_file():
                continue
            logical_id = str(img.get("logical_id") or f"s{idx+1}_img{j+1}")[:40]
            left = Inches(pic_left + j * (pic_width + 0.25))
            try:
                pic = slide.shapes.add_picture(
                    str(img_path),
                    left,
                    Inches(4.6),
                    width=Inches(pic_width),
                )
                media_count += 1
                logical_map[logical_id] = {
                    "slide": idx + 2,
                    "shape_id": pic.shape_id,
                    "name": logical_id,
                }
            except Exception:
                continue

        if notes:
            try:
                slide.notes_slide.notes_text_frame.text = notes
            except Exception:
                pass

    bio = io.BytesIO()
    prs.save(bio)
    output_path.write_bytes(bio.getvalue())
    return {
        "slide_count": len(slides_in) + 1,
        "media_count": media_count,
        "logical_map": logical_map,
    }


def copy_template_base(template_path: Path, output_path: Path) -> Dict[str, Any]:
    output_path = Path(output_path)
    template_path = Path(template_path)
    if not template_path.is_file():
        raise FileNotFoundError(f"模板不存在：{template_path}")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(template_path, output_path)
    try:
        from pptx import Presentation

        prs = Presentation(str(output_path))
        slide_count = len(prs.slides)
    except Exception:
        slide_count = 1
    return {"slide_count": slide_count, "media_count": 0, "logical_map": {}}
