"""PPT 全量读取与 PPT 生成员工：检测、规则、兜底 convert 与包体验证。"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any, Dict, List, Tuple

from modstore_server.ppt_convert_templates import (
    render_ppt_generate_convert_module as render_ppt_generate_convert_module,
    render_ppt_read_convert_module as render_ppt_read_convert_module,
)

PPT_DOC_KEYWORDS = (
    ".pptx",
    ".ppt",
    "pptx",
    "ppt",
    "powerpoint",
    "演示文稿",
    "幻灯片",
    "presentation",
)
PPT_READ_ACTION_KEYWORDS = (
    "读取",
    "读出",
    "全量",
    "读入",
    "read",
    "load",
    "提取",
    "解析",
    "演讲",
    "备注",
    "备注稿",
    "大纲",
    "vlm",
    "识图",
)
PPT_GENERATE_ACTION_KEYWORDS = (
    "生成",
    "写入",
    "写 ppt",
    "写ppt",
    "输出",
    "改写",
    "write",
    "generate",
    "json",
    "结构化",
    "中介",
    "制作",
)
PPT_READ_EXCLUDE = (
    "生成 ppt",
    "写 ppt",
    "写ppt",
    "json 中介生成",
    "写出 output.pptx",
)
PPT_GENERATE_EXCLUDE = (
    "仅读取",
    "只读",
    "原样",
    "不要生成",
    "read only",
    "全量读取",
    "演讲备注生成",
)

PPT_READ_OUTPUT_FIELDS = (
    "title",
    "slide_count",
    "outline",
    "slides",
    "images",
    "notes_generated",
    "source",
)
PPT_GENERATE_OUTPUT_FIELDS = (
    "title",
    "slides",
    "slide_count",
    "stats",
    "metadata",
)

IMAGE_CATEGORY_DIRS = ("figures", "photos", "diagrams", "icons", "uncategorized")

SPEAKER_NOTES_PROMPT = "为这份PPT生成每页的演讲备注"


def _brief_lower(brief: str) -> str:
    return (brief or "").lower()


def _has_ppt_doc_signal(bl: str) -> bool:
    return any(k in bl for k in PPT_DOC_KEYWORDS)


def is_ppt_generate(brief: str) -> bool:
    """上传 JSON（presentation_full 同 schema）→ 写出 output.pptx。"""
    bl = _brief_lower(brief)
    if not _has_ppt_doc_signal(bl):
        return False
    if any(k in bl for k in PPT_GENERATE_EXCLUDE) and not any(
        k in bl for k in PPT_GENERATE_ACTION_KEYWORDS
    ):
        return False
    if any(k in bl for k in PPT_READ_EXCLUDE):
        return False
    return any(k in bl for k in PPT_GENERATE_ACTION_KEYWORDS)


def is_ppt_full_read(brief: str) -> bool:
    """PPT 全量读取：大纲、每页正文、图片 VLM、演讲备注。"""
    if is_ppt_generate(brief):
        return False
    bl = _brief_lower(brief)
    if not _has_ppt_doc_signal(bl):
        return False
    return any(k in bl for k in PPT_READ_ACTION_KEYWORDS)


def ppt_read_structured_spec(brief: str) -> Dict[str, Any]:
    return {
        "domain": "文档处理 / PPT 全量读取",
        "goal": (brief or "").strip().splitlines()[0][:200] or "上传 PPT，全量解析并生成演讲备注",
        "input": "用户上传的 .pptx 文件",
        "output": "outputs/presentation_full.json + speaker_notes.md + outputs/images/",
        "output_schema": {
            "fields": list(PPT_READ_OUTPUT_FIELDS),
            "json_file": "outputs/presentation_full.json",
            "meta_file": "outputs/presentation_meta.json",
            "notes_file": "outputs/speaker_notes.md",
            "images_index": "outputs/images_index.json",
            "images_dir": "outputs/images/",
        },
        "constraints": [
            "幻灯片正文必须来自 python-pptx 真实解析，禁止 LLM 编造正文",
            "图片须导出并由 VLM（可用时）生成 sidecar",
            f"演讲备注由 LLM 基于真实页内容生成，提示词：{SPEAKER_NOTES_PROMPT}",
            "handlers 必须为 direct_python",
        ],
        "suggested_capabilities": ["ppt.parse", "ppt.notes_generate", "vision.vlm"],
        "suggested_handlers": ["direct_python"],
    }


def ppt_generate_structured_spec(brief: str) -> Dict[str, Any]:
    return {
        "domain": "文档处理 / PPT 生成",
        "goal": (brief or "").strip().splitlines()[0][:200] or "LLM 编排 + OOXML 写出 output.pptx",
        "input": "presentation_full v2 JSON / user_query / .txt；可选 template .pptx（execute-file template_file）",
        "output": "outputs/output.pptx + outputs/ppt_edit_plan.json",
        "output_schema": {
            "fields": list(PPT_GENERATE_OUTPUT_FIELDS),
            "pptx_file": "outputs/output.pptx",
            "plan_file": "outputs/ppt_edit_plan.json",
        },
        "constraints": [
            "compose-first：无模板时从零合成多页 pptx",
            "enhance：复制 template 后按 ppt_edit_plan 注入动画（OOXML）",
            "禁止仅输出纯文字幻灯片冒充带动效/带图作业",
            "作业跑马灯可走 homework_marquee 确定性配方",
        ],
        "suggested_capabilities": ["ppt.write", "ppt.ooxml", "data.json_read", "llm.plan"],
        "suggested_handlers": ["direct_python", "agent"],
    }


def build_ppt_read_rule_spec(brief: str) -> Dict[str, Any]:
    return {
        "brief": brief,
        "mode": "direct_python_file_transform",
        "accepted_extensions": [".pptx"],
        "default_action": "convert",
        "default_output_relpath": "outputs/presentation_full.json",
        "default_meta_relpath": "outputs/presentation_meta.json",
        "default_images_dir": "outputs/images",
        "default_notes_relpath": "outputs/speaker_notes.md",
        "runtime_kind": "ppt_full_read",
        "speaker_notes_prompt": SPEAKER_NOTES_PROMPT,
        "output_schema": list(PPT_READ_OUTPUT_FIELDS),
        "requirements": [
            'Use direct_python only; handlers must be ["direct_python"].',
            "Parse pptx with python-pptx; never use LLM for slide body text.",
            "Export embedded images to outputs/images/<category>/.",
            "When ctx.call_llm supports vision, describe each image to .vlm.json sidecar.",
            f"Generate per-slide speaker notes via ctx.call_llm text with prompt: {SPEAKER_NOTES_PROMPT}",
            "Write presentation_full.json, presentation_meta.json, speaker_notes.md, images_index.json.",
            "Return {ok, summary, items, warnings, error, meta}.",
        ],
    }


def build_ppt_generate_rule_spec(brief: str) -> Dict[str, Any]:
    bl = _brief_lower(brief)
    wants_polish = any(k in bl for k in ("润色", "改写", "polish", "rewrite"))
    return {
        "brief": brief,
        "mode": "direct_python_file_transform",
        "accepted_extensions": [".json", ".txt"],
        "default_action": "convert",
        "default_output_relpath": "outputs/output.pptx",
        "runtime_kind": "ppt_generate",
        "optional_llm_polish": wants_polish,
        "output_schema": list(PPT_GENERATE_OUTPUT_FIELDS),
        "requirements": [
            'handlers must include "direct_python"; may include "agent" for optional polish.',
            "Run modstore_server.ppt_generate_pipeline.run_ppt_generate: route → plan → compose/enhance → OOXML.",
            "When template_path is .pptx, copy template; else compose deck from plan slides.",
            "LLM planner outputs ppt_edit_plan.json; executor applies inject_timing presets.",
            "Never fabricate slides when input is empty; fallback to text-only only with warnings.",
            "Return {ok, summary, items, warnings, error, meta}.",
        ],
    }


def validate_ppt_read_backend(pack_dir: Path) -> Tuple[List[str], List[str]]:
    return _validate_ppt_backend(
        pack_dir,
        runtime_kind="ppt_full_read",
        required_tokens=("presentation_full", "speaker_notes", ".pptx"),
    )


def validate_ppt_generate_backend(pack_dir: Path) -> Tuple[List[str], List[str]]:
    return _validate_ppt_backend(
        pack_dir,
        runtime_kind="ppt_generate",
        required_tokens=("output.pptx", "presentation"),
    )


def _validate_ppt_backend(
    pack_dir: Path,
    *,
    runtime_kind: str,
    required_tokens: tuple[str, ...],
) -> Tuple[List[str], List[str]]:
    errors: List[str] = []
    warnings: List[str] = []
    backend = pack_dir / "backend"
    if not backend.is_dir():
        errors.append("缺少 backend 目录")
        return errors, warnings

    py_blob = ""
    has_convert = False
    for py_path in backend.rglob("*.py"):
        try:
            text = py_path.read_text(encoding="utf-8", errors="ignore")
            py_blob += text.lower()
            if "def convert_file" in text and "vendor" in str(py_path).lower():
                has_convert = True
        except OSError:
            pass

    mf_path = pack_dir / "manifest.json"
    handlers: List[str] = []
    if mf_path.is_file():
        try:
            from modstore_server.employee_asset_pipeline import manifest_actions_handlers

            mf = json.loads(mf_path.read_text(encoding="utf-8"))
            handlers = manifest_actions_handlers(mf)
        except (json.JSONDecodeError, OSError):
            warnings.append("manifest.json 无法解析")

    if handlers and "direct_python" not in handlers:
        errors.append(f"{runtime_kind} 员工 handlers 必须包含 direct_python")
    if not has_convert:
        errors.append("backend/vendor 中缺少 convert_file 实现")
    if "pptx" not in py_blob and "presentation" not in py_blob:
        warnings.append("未发现 PPT 解析相关代码")
    for tok in required_tokens:
        if tok.lower() not in py_blob:
            warnings.append(f"convert 模块可能未覆盖：{tok}")

    return errors, warnings


def minimal_pptx_fixture_bytes() -> bytes:
    try:
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "PPT 冒烟测试"
        body = slide.placeholders[1]
        body.text_frame.text = "要点一\n要点二"
        import io

        bio = io.BytesIO()
        prs.save(bio)
        return bio.getvalue()
    except Exception:
        return b""


def ppt_read_orchestration_plan(brief: str, payload: Dict[str, Any]) -> Dict[str, Any]:
    from modstore_server.employee_brief_utils import compact_routing_brief

    checklist = payload.get("execution_checklist")
    checklist_text = (
        "\n".join(f"- {x}" for x in checklist if isinstance(x, str))
        if isinstance(checklist, list)
        else ""
    )
    clean = compact_routing_brief(brief, max_len=400) or (brief or "").strip()
    merged = "\n".join(x for x in [clean, checklist_text] if x).strip()
    short = "PPT 全量读取员"
    script_brief = (
        f"{merged or clean}\n\n"
        "请生成 Python：读取 inputs/ 中 .pptx，解析每页正文与大纲，"
        "导出图片并 VLM 描述，按「为这份PPT生成每页的演讲备注」生成 notes_generated，"
        "写入 outputs/presentation_full.json 与 speaker_notes.md。"
    )
    return {
        "employee_name": short,
        "employee_brief": (
            f"{merged or clean}\n\n"
            "员工必须使用 direct_python：正文仅来自 pptx 解析；图片走 VLM；演讲备注基于真实内容生成。"
        ),
        "script_workflow_name": f"{short} 脚本工作流",
        "script_brief": script_brief,
        "script_runtime_notes": "只能读 inputs/、写 outputs/；VLM/备注通过 ctx.call_llm。",
        "workflow_name": str(payload.get("employee_workflow_name") or short).strip() or short,
        "workflow_brief": f"{merged or clean}\n\nSkill：上传 pptx → JSON 中介 + 演讲备注 + 图片 VLM。",
        "acceptance": [
            "handlers 为 direct_python",
            "presentation_full.json 含 slides/outline",
            "speaker_notes.md 含每页备注",
        ],
    }


def ppt_generate_orchestration_plan(brief: str, payload: Dict[str, Any]) -> Dict[str, Any]:
    from modstore_server.employee_brief_utils import compact_routing_brief

    checklist = payload.get("execution_checklist")
    checklist_text = (
        "\n".join(f"- {x}" for x in checklist if isinstance(x, str))
        if isinstance(checklist, list)
        else ""
    )
    clean = compact_routing_brief(brief, max_len=400) or (brief or "").strip()
    merged = "\n".join(x for x in [clean, checklist_text] if x).strip()
    short = "PPT 生成员"
    script_brief = (
        f"{merged or clean}\n\n"
        "请生成 Python：读取 inputs/ .json（presentation_full schema）→ 写出 outputs/output.pptx。"
    )
    return {
        "employee_name": short,
        "employee_brief": (
            f"{merged or clean}\n\n" "JSON 为中介；direct_python 写 pptx；禁止无输入编造幻灯片。"
        ),
        "script_workflow_name": f"{short} 脚本工作流",
        "script_brief": script_brief,
        "script_runtime_notes": "direct_python 读 JSON 写 pptx。",
        "workflow_name": str(payload.get("employee_workflow_name") or short).strip() or short,
        "workflow_brief": f"{merged or clean}\n\nSkill：JSON → output.pptx。",
        "acceptance": [
            "输出 output.pptx",
            "handlers 含 direct_python",
        ],
    }


def resolve_ppt_orchestration_plan(brief: str, payload: Dict[str, Any]) -> Dict[str, Any]:
    if is_ppt_generate(brief):
        return ppt_generate_orchestration_plan(brief, payload)
    return ppt_read_orchestration_plan(brief, payload)
