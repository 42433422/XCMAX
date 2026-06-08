"""Enrich presentation_full to v2 with per-shape geometry and timing hints."""

from __future__ import annotations

import zipfile
from pathlib import Path
from typing import Any, Dict, List


def _deck_has_timing(pptx_path: Path) -> bool:
    try:
        with zipfile.ZipFile(pptx_path, "r") as z:
            for name in z.namelist():
                if name.startswith("ppt/slides/slide") and b"timing>" in z.read(name):
                    return True
    except Exception:
        return False
    return False


def enrich_presentation_v2(presentation: Dict[str, Any], src_path: Path) -> Dict[str, Any]:
    src_path = Path(src_path)
    out = dict(presentation or {})
    out["schema_version"] = 2
    out["has_timing"] = _deck_has_timing(src_path)
    out["animation_hints"] = ["click_sequence"] if out["has_timing"] else []

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return out

    prs = Presentation(str(src_path))
    slides_in = out.get("slides") if isinstance(out.get("slides"), list) else []
    slide_by_index = {
        int(s.get("index") or i + 1): s for i, s in enumerate(slides_in) if isinstance(s, dict)
    }

    for slide_idx, slide in enumerate(prs.slides, start=1):
        shapes_out: List[Dict[str, Any]] = []
        img_seq = 0
        for shape in slide.shapes:
            sid = int(getattr(shape, "shape_id", 0) or 0)
            name = str(getattr(shape, "name", "") or "")
            left = int(getattr(shape, "left", 0) or 0)
            top = int(getattr(shape, "top", 0) or 0)
            width = int(getattr(shape, "width", 0) or 0)
            height = int(getattr(shape, "height", 0) or 0)
            stype = "unknown"
            embed_rid = ""
            image_relpath = ""
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                stype = "picture"
                img_seq += 1
                logical_id = f"s{slide_idx}_img{img_seq}"
                try:
                    image_relpath = f"figures/{logical_id}.png"
                except Exception:
                    pass
            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or getattr(
                shape, "has_text_frame", False
            ):
                stype = "text"
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                stype = "group"
            shapes_out.append(
                {
                    "logical_id": name or f"shape_{sid}",
                    "shape_id": sid,
                    "type": stype,
                    "name": name,
                    "x": left,
                    "y": top,
                    "cx": width,
                    "cy": height,
                    "embed_rid": embed_rid,
                    "image_relpath": image_relpath,
                }
            )
        entry = slide_by_index.get(slide_idx)
        if entry is None:
            entry = {"index": slide_idx, "title": f"第 {slide_idx} 页", "bullets": [], "images": []}
            slides_in.append(entry)
        entry["shapes"] = shapes_out

    out["slides"] = slides_in
    out["slide_count"] = len(prs.slides)
    return out
